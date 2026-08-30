"""Format coverage tests: pptx + legacy xls parsers and detectors."""

from __future__ import annotations

from aiitg.core.registry import default_format_registry
from tests.fixtures import builders


class TestPptx:
    def test_benign_parses(self, scan_file, tmp_path):
        f = builders.build_pptx_benign(tmp_path / "ok.pptx")
        doc = default_format_registry().parse(f)
        assert doc.kind == "pptx"
        assert "benign slide" in doc.all_text

    def test_white_text_detected(self, scan_file, tmp_path):
        f = builders.build_pptx_white_text(tmp_path / "white.pptx")
        report = scan_file(f)
        hits = [ev for ev in report.evidence if ev.detector_id == "DET-002"]
        assert len(hits) >= 1
        assert "white" in hits[0].description.lower()

    def test_notes_present(self, scan_file, tmp_path):
        f = builders.build_pptx_notes(tmp_path / "notes.pptx")
        doc = default_format_registry().parse(f)
        assert "SECRET INSTRUCTIONS" in doc.all_text

    def test_benign_no_evidence(self, scan_file, tmp_path):
        f = builders.build_pptx_benign(tmp_path / "ok.pptx")
        report = scan_file(f)
        assert len(report.evidence) == 0

    def test_pptx_zerowidth(self, scan_file, tmp_path):
        # pptx with zero-width char in a text box
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        box.text = "Visible\u200bINSTRUCTION"
        f = tmp_path / "zw.pptx"
        prs.save(str(f))
        report = scan_file(f)
        hits = [ev for ev in report.evidence if ev.detector_id == "DET-001"]
        assert len(hits) >= 1


class TestXls:
    def test_benign_parses(self, scan_file, tmp_path):
        f = builders.build_xls_benign(tmp_path / "ok.xls")
        doc = default_format_registry().parse(f)
        assert doc.kind == "xls"
        assert "hello" in doc.all_text

    def test_benign_no_evidence(self, scan_file, tmp_path):
        f = builders.build_xls_benign(tmp_path / "ok.xls")
        report = scan_file(f)
        assert len(report.evidence) == 0

    def test_xls_zerowidth(self, scan_file, tmp_path):
        import xlwt

        wb = xlwt.Workbook()
        ws = wb.add_sheet("S")
        ws.write(0, 0, "normal")
        ws.write(0, 1, "secret\u200binstruction")
        f = tmp_path / "zw.xls"
        wb.save(str(f))
        report = scan_file(f)
        hits = [ev for ev in report.evidence if ev.detector_id == "DET-001"]
        assert len(hits) >= 1
