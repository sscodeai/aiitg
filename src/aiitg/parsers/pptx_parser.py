"""pptx parser: python-pptx → unified text model (slides, shapes, notes).

PowerPoint files are a common LLM ingestion target (slide decks → summaries,
pitch-deck analysis). Hidden content in slides:
- off-slide / tiny-font text boxes
- white text on white background
- notes pages (speaker notes are machine-readable but often ignored by humans)
- zero-width chars in any text frame
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation

from aiitg.core.document import ParsedDocument, TextParagraph, TextRun
from aiitg.parsers.base import register_parser


def _normalize_color(val: str | None) -> str | None:
    if not val:
        return None
    v = val.strip().lstrip("#").upper()
    if len(v) == 6 and all(c in "0123456789ABCDEF" for c in v):
        return v
    if len(v) == 3 and all(c in "0123456789ABCDEF" for c in v):
        return "".join(c * 2 for c in v)
    return None


def _run_from_font(text: str, font, is_notes: bool = False) -> TextRun:
    """Build a TextRun from a pptx font object."""
    size = None
    color = None
    try:
        if font.size is not None:
            size = float(font.size.pt)
    except Exception:  # noqa: BLE001
        size = None
    try:
        if font.color and font.color.type is not None and font.color.rgb is not None:
            color = _normalize_color(str(font.color.rgb))
    except Exception:  # noqa: BLE001
        color = None
    return TextRun(text=text, font_size=size, color=color, is_header_footer=is_notes)


@register_parser("pptx", ("pptx",), sniff=None)
class PptxParser:
    """Parse .pptx into :class:`ParsedDocument`."""

    def parse(self, path: Path) -> ParsedDocument:
        prs = Presentation(str(path))
        doc_model = ParsedDocument(kind="pptx", source_path=str(path))
        slide_num = 0

        for slide in prs.slides:
            slide_num += 1
            for shape in slide.shapes:
                text_frame = getattr(shape, "text_frame", None)
                if text_frame is None:
                    continue
                for para in text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs)
                    if not text.strip():
                        continue
                    runs: list[TextRun] = []
                    for r in para.runs:
                        runs.append(_run_from_font(r.text, r.font))
                    doc_model.paragraphs.append(
                        TextParagraph(
                            text=text,
                            index=slide_num,
                            runs=runs,
                            raw={"slide": slide_num, "shape": shape.shape_id, "element": f"slide{slide_num}"},
                        )
                    )

            # notes pages (speaker notes — machine-readable, human-invisible)
            if slide.has_notes_slide and slide.notes_slide is not None:
                notes_frame = slide.notes_slide.notes_text_frame
                if notes_frame is not None and notes_frame.text.strip():
                    notes_text = notes_frame.text
                    doc_model.paragraphs.append(
                        TextParagraph(
                            text=notes_text,
                            index=slide_num,
                            runs=[TextRun(text=notes_text, is_header_footer=True, raw={"element": "notes"})],
                            raw={"slide": slide_num, "element": "notes"},
                        )
                    )

        # core properties
        try:
            props = prs.core_properties
            doc_model.metadata = {
                "author": props.author,
                "title": props.title,
                "subject": props.subject,
                "keywords": props.keywords,
                "comments": props.comments,
            }
        except Exception:  # noqa: BLE001
            doc_model.metadata = {}

        # raw OOXML parts (for raw-node detector)
        try:
            from aiitg.parsers.ooxml_raw import OOXMLRawReader

            with OOXMLRawReader(path) as reader:
                doc_model.ooxml_parts = {name: reader.part_bytes(name) for name in reader.part_names()}
        except Exception as exc:  # noqa: BLE001
            doc_model.warnings.append(f"failed to read raw OOXML parts: {exc}")

        return doc_model
